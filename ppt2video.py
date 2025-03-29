import tkinter as tk
from tkinter import ttk
from tkinter import filedialog, messagebox
from pptx import Presentation
import os
import subprocess
from PIL import Image
import json
import time
import comtypes.client
import shutil


class PPTToVideo:
    def __init__(self):
        self.window = tk.Tk()
        self.window.title("PPT转视频工具 V1.0")
        self.window.geometry("800x700")
        self.window.resizable(False, False)  # 禁止调整窗口大小

        # 设置主题颜色
        self.orange_color = "#F25022"
        self.window.configure(bg="#FFFFFF")

        # 文件路径变量
        self.ppt_path = tk.StringVar()
        self.bgm_path = tk.StringVar()

        # 参数变量及默认值设置
        self.slide_duration = tk.StringVar(value="5")
        self.transition_duration = tk.StringVar(value="1")
        self.video_quality = tk.StringVar(value="高质量")
        self.bgm_volume = tk.StringVar(value="1.0")
        self.auto_next = tk.BooleanVar(value=True)
        self.save_text = tk.BooleanVar(value=True)
        self.resolution = tk.StringVar(value="自动")

        self.transition_effect = tk.StringVar(value="无")
        self.transition_effects = {
            "无": "",
            "淡入淡出": "fade",
            "向左滑动": "slideleft",
            "向右滑动": "slideright",
            "向上滑动": "slideup",
            "向下滑动": "slidedown",
            "随机效果": "random"
        }

        # 日志相关
        self.log_file = "ffmpeg_log.txt"
        self.log_window = None

        # 配置文件路径
        self.config_path = "config.json"
        self.load_config()

        # 确保temp目录存在
        self.temp_dir = os.path.join(os.getcwd(), "TEMP")
        os.makedirs(self.temp_dir, exist_ok=True)

        self.create_widgets()

    def create_widgets(self):
        # 创建自定义样式
        style = ttk.Style()
        style.configure("Orange.TButton",
                        padding=5,
                        relief="flat",
                        background=self.orange_color)
        style.configure("Section.TLabelframe",
                        padding=10,
                        relief="solid",
                        borderwidth=1)

        # 主框架
        main_frame = ttk.Frame(self.window, padding=20)
        main_frame.pack(fill='both', expand=True)

        # 1. 文件选择区域
        file_frame = ttk.LabelFrame(main_frame,
                                    text=" 文件选择 ",
                                    style="Section.TLabelframe")
        file_frame.pack(fill='x', pady=(0, 15))

        # PPT文件选择行
        ppt_frame = ttk.Frame(file_frame)
        ppt_frame.pack(fill='x', pady=5)
        ttk.Label(ppt_frame, text="PPT文件：", width=10).pack(side='left')
        ttk.Entry(ppt_frame, textvariable=self.ppt_path).pack(side='left',
                                                              fill='x',
                                                              expand=True,
                                                              padx=5)
        ttk.Button(ppt_frame,
                   text="浏览",
                   command=self.select_ppt,
                   style="Orange.TButton").pack(side='left')

        # 背景音乐选择行
        bgm_frame = ttk.Frame(file_frame)
        bgm_frame.pack(fill='x', pady=5)
        ttk.Label(bgm_frame, text="背景音乐：", width=10).pack(side='left')
        ttk.Entry(bgm_frame, textvariable=self.bgm_path).pack(side='left',
                                                              fill='x',
                                                              expand=True,
                                                              padx=5)
        ttk.Button(bgm_frame,
                   text="浏览",
                   command=self.select_bgm,
                   style="Orange.TButton").pack(side='left')

        # 2. 参数设置区域
        params_frame = ttk.LabelFrame(main_frame,
                                      text=" 参数设置 ",
                                      style="Section.TLabelframe")
        params_frame.pack(fill='x', pady=(0, 15))

        # 创建网格布局
        for i in range(4):
            params_frame.grid_columnconfigure(i, weight=1, pad=10)

        # 第一行参数
        ttk.Label(params_frame, text="每页停留时间(秒)：").grid(row=0,
                                                        column=0,
                                                        sticky='e',
                                                        pady=5)
        ttk.Entry(params_frame,
                  textvariable=self.slide_duration,
                  width=10).grid(row=0, column=1, sticky='w')

        ttk.Label(params_frame, text="转场时间(秒)：").grid(row=0,
                                                      column=2,
                                                      sticky='e')
        ttk.Entry(params_frame,
                  textvariable=self.transition_duration,
                  width=10).grid(row=0, column=3, sticky='w')

        # 第二行参数
        ttk.Label(params_frame, text="视频质量：").grid(row=1,
                                                   column=0,
                                                   sticky='e',
                                                   pady=5)
        ttk.Combobox(params_frame,
                     textvariable=self.video_quality,
                     values=["低质量", "中等质量", "高质量"],
                     state='readonly',
                     width=10).grid(row=1, column=1, sticky='w')

        ttk.Label(params_frame, text="分辨率：").grid(row=1,
                                                  column=2,
                                                  sticky='e')
        ttk.Combobox(params_frame,
                     textvariable=self.resolution,
                     values=["自动", "1280x720", "1920x1080", "2560x1440"],
                     state='readonly',
                     width=10).grid(row=1, column=3, sticky='w')


        # 第三行参数
        ttk.Label(params_frame, text="背景音量(0-1)：").grid(row=2,
                                                        column=0,
                                                        sticky='e',
                                                        pady=5)
        ttk.Entry(params_frame,
                  textvariable=self.bgm_volume,
                  width=10).grid(row=2, column=1, sticky='w')

        ttk.Label(params_frame, text="转场效果：").grid(row=2, column=2, sticky='e')
        ttk.Combobox(params_frame,
                     textvariable=self.transition_effect,
                     values=list(self.transition_effects.keys()),
                     state='readonly',
                     width=10).grid(row=2, column=3, sticky='w')

        # 选项区域
        options_frame = ttk.Frame(params_frame)
        options_frame.grid(row=3, column=0, columnspan=4, pady=10)
        ttk.Checkbutton(options_frame,
                        text="自动翻页",
                        variable=self.auto_next).pack(side='left', padx=20)
        ttk.Checkbutton(options_frame,
                        text="保存提取文本",
                        variable=self.save_text).pack(side='left', padx=20)

        # 3. 操作按钮区域
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(pady=20)

        buttons = [
            ("保存设置", self.save_config),
            ("编辑文案", self.edit_text),
            ("查看日志", self.show_log),
            ("开始转换", self.convert),
            ("另存视频", self.save_video_as)
        ]

        for text, command in buttons:
            ttk.Button(button_frame,
                       text=text,
                       command=command,
                       style="Orange.TButton").pack(side='left', padx=10)

    def convert_ppt_to_images(self, ppt_path, output_dir):
        try:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1

            ppt = powerpoint.Presentations.Open(ppt_path)
            total_slides = ppt.Slides.Count

            self.show_progress_window(total_slides)

            if self.resolution.get() == "自动":
                width, height = 0, 0
            else:
                width, height = map(int, self.resolution.get().split('x'))

            for i in range(1, total_slides + 1):
                slide = ppt.Slides(i)
                image_path = os.path.join(output_dir, f"slide_{i}.png")

                if width and height:
                    slide.Export(image_path, "PNG", width, height)
                else:
                    slide.Export(image_path, "PNG")

                while not os.path.exists(image_path):
                    time.sleep(0.1)

                with Image.open(image_path) as img:
                    img.save(image_path, "PNG", optimize=True)

                self.update_progress(i, total_slides)

            ppt.Close()
            powerpoint.Quit()
            self.progress_window.destroy()
            return True

        except Exception as e:
            if hasattr(self, 'progress_window'):
                self.progress_window.destroy()
            messagebox.showerror("错误", f"转换PPT到图片失败：{str(e)}")
            return False
        finally:
            try:
                powerpoint.Quit()
            except:
                pass

    def convert(self):
        if not self.ppt_path.get() or not self.bgm_path.get():
            messagebox.showerror("错误", "请选择PPT文件和背景音乐！")
            return

        try:
            # 清空并创建临时目录
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
            os.makedirs(self.temp_dir)

            # 清空日志文件
            with open(self.log_file, 'w') as f:
                f.write("")

            # 提取文本并保存
            if self.save_text.get():
                self.extract_text_from_ppt(self.ppt_path.get())

            # 转换PPT到图片
            if not self.convert_ppt_to_images(
                    os.path.abspath(self.ppt_path.get()),
                    self.temp_dir
            ):
                raise Exception("PPT转图片失败")

            # 获取尺寸信息
            if self.resolution.get() == "自动":
                width, height = self.get_max_slide_dimensions(self.temp_dir)
            else:
                width, height = map(int, self.resolution.get().split('x'))

            # 准备转场效果
            transition_name = self.transition_effects[self.transition_effect.get()]
            if transition_name == "random":
                import random
                effects = list(self.transition_effects.values())[1:-1]  # 排除"无"和"随机"
                transition_name = random.choice(effects)

            # 创建图片列表文件
            image_files = [f for f in os.listdir(self.temp_dir) if f.endswith('.png')]
            image_files.sort(key=lambda x: int(x.split('_')[1].split('.')[0]))
            total_duration = len(image_files) * float(self.slide_duration.get())

            with open(os.path.join(self.temp_dir, "input.txt"), "w", encoding="utf-8") as f:
                for i, image_file in enumerate(image_files):
                    image_path = os.path.abspath(os.path.join(self.temp_dir, image_file))
                    image_path = image_path.replace('\\', '/').encode('utf-8').decode('utf-8')
                    f.write(f"file '{image_path}'\n")
                    if i < len(image_files) - 1:
                        f.write(f"duration {self.slide_duration.get()}\n")
                    else:
                        f.write("duration 1\n")  # 最后一帧持续1秒

            # 构建filter_complex
            def build_filter_complex():
                filters = []
                base_filter = f"scale={width}:{height}:force_original_aspect_ratio=decrease,pad={width}:{height}:(ow-iw)/2:(oh-ih)/2:black"

                if transition_name:
                    filters.append(f"[0:v]{base_filter},split[v1][v2];")
                    filters.append(
                        f"[v1][v2]xfade=transition={transition_name}:duration={self.transition_duration.get()}[vout]")
                else:
                    filters.append(f"[0:v]{base_filter}[vout]")

                # 音频处理
                filters.append(f"[1:a]volume={self.bgm_volume.get()},afade=t=out:st={total_duration - 3}:d=3[aout]")

                return ";".join(filters)

            # 构建FFmpeg命令
            cmd = [
                      "ffmpeg", "-y",
                      "-f", "concat",
                      "-safe", "0",
                      "-i", os.path.join(self.temp_dir, "input.txt"),
                      "-i", self.bgm_path.get(),
                      "-filter_complex", build_filter_complex(),
                      "-map", "[vout]",
                      "-map", "[aout]",
                      "-c:v", "libx264"
                  ] + self.get_ffmpeg_quality_params() + [
                      "-pix_fmt", "yuv420p",
                      "-shortest",
                      os.path.join(self.temp_dir, "output.mp4")
                  ]

            # 显示进度窗口
            progress_window = tk.Toplevel(self.window)
            progress_window.title("转换进度")
            progress_window.geometry("300x150")
            progress_window.transient(self.window)
            progress_window.grab_set()

            progress_label = ttk.Label(progress_window, text="正在转换视频...")
            progress_label.pack(pady=20)

            progress_bar = ttk.Progressbar(progress_window, mode='indeterminate')
            progress_bar.pack(pady=20)
            progress_bar.start()

            print(cmd)

            # 运行FFmpeg并记录日志
            with open(self.log_file, 'a', encoding='utf-8') as log:
                process = subprocess.Popen(
                    cmd,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    universal_newlines=True,
                    encoding='utf-8'
                )

                while True:
                    output = process.stderr.readline()
                    if output == '' and process.poll() is not None:
                        break
                    if output:
                        log.write(output)
                        log.flush()
                    progress_window.update()

                progress_bar.stop()
                progress_window.destroy()

                if process.returncode != 0:
                    raise Exception("FFmpeg转换失败")

            messagebox.showinfo("成功", "转换完成！请使用'另存视频'功能保存到指定位置。")

        except Exception as e:
            if 'progress_window' in locals():
                progress_window.destroy()
            messagebox.showerror("错误", f"转换失败：{str(e)}")
            # 记录详细错误信息到日志
            with open(self.log_file, 'a', encoding='utf-8') as log:
                log.write(f"\nError occurred: {str(e)}\n")
        finally:
            # 确保进度窗口被关闭
            if 'progress_window' in locals():
                try:
                    progress_window.destroy()
                except:
                    pass

    def extract_text_from_ppt(self, ppt_path):
        output_file = "ppt_content.txt"
        try:
            prs = Presentation(ppt_path)
            text_extracted = False

            with open(output_file, 'w', encoding='utf-8') as f:
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    f.write(f"=== 第{i}页 ===\n")

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                            text_extracted = True

                    if slide_text:
                        f.write('\n'.join(slide_text) + '\n')
                    f.write('\n')

            if not text_extracted:
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write("未检测到PPT中的文本内容\n")

        except Exception as e:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"提取文本时发生错误: {str(e)}\n")

    def show_progress_window(self, total_slides):
        self.progress_window = tk.Toplevel(self.window)
        self.progress_window.title("转换进度")
        self.progress_window.geometry("300x150")
        self.progress_window.transient(self.window)
        self.progress_window.grab_set()

        # 居中显示
        self.progress_window.update_idletasks()
        width = self.progress_window.winfo_width()
        height = self.progress_window.winfo_height()
        x = (self.progress_window.winfo_screenwidth() // 2) - (width // 2)
        y = (self.progress_window.winfo_screenheight() // 2) - (height // 2)
        self.progress_window.geometry(f'{width}x{height}+{x}+{y}')

        # 进度显示组件
        self.progress_label = ttk.Label(self.progress_window, text="正在转换PPT为图片...")
        self.progress_label.pack(pady=10)

        self.progress_bar = ttk.Progressbar(self.progress_window, length=200, mode='determinate')
        self.progress_bar.pack(pady=10)

        self.slide_count_label = ttk.Label(self.progress_window, text=f"0/{total_slides} 页")
        self.slide_count_label.pack(pady=10)

        self.progress_window.update()

    def update_progress(self, current, total):
        progress = (current / total) * 100
        self.progress_bar['value'] = progress
        self.slide_count_label.config(text=f"{current}/{total} 页")
        self.progress_window.update()

    def get_max_slide_dimensions(self, temp_dir):
        max_width = 0
        max_height = 0
        for file in os.listdir(temp_dir):
            if file.startswith("slide_") and file.endswith(".png"):
                with Image.open(os.path.join(temp_dir, file)) as img:
                    width, height = img.size
                    max_width = max(max_width, width)
                    max_height = max(max_height, height)

        # 确保尺寸为偶数
        max_width = max_width + (max_width % 2)
        max_height = max_height + (max_height % 2)

        return max_width, max_height

    def edit_text(self):
        if not os.path.exists("ppt_content.txt"):
            messagebox.showerror("错误", "未找到文本文件，请先转换PPT")
            return

        # 使用系统默认编辑器打开文本文件
        if os.name == 'nt':  # Windows
            os.startfile("ppt_content.txt")
        else:  # Linux/Mac
            subprocess.call(('xdg-open', "ppt_content.txt"))

    def show_log(self):
        if self.log_window is None or not self.log_window.winfo_exists():
            self.log_window = tk.Toplevel(self.window)
            self.log_window.title("转换日志")
            self.log_window.geometry("600x400")

            text_widget = tk.Text(self.log_window)
            text_widget.pack(fill='both', expand=True)

            scrollbar = ttk.Scrollbar(self.log_window, orient='vertical', command=text_widget.yview)
            scrollbar.pack(side='right', fill='y')
            text_widget.configure(yscrollcommand=scrollbar.set)

            def update_log():
                if os.path.exists(self.log_file):
                    try:
                        with open(self.log_file, 'r', encoding='utf-8') as f:
                            text_widget.delete(1.0, tk.END)
                            text_widget.insert(tk.END, f.read())
                    except Exception as e:
                        text_widget.delete(1.0, tk.END)
                        text_widget.insert(tk.END, f"读取日志出错: {str(e)}")
                self.log_window.after(1000, update_log)

            update_log()

    def save_video_as(self):
        output_video = os.path.join(self.temp_dir, "output.mp4")
        if not os.path.exists(output_video):
            messagebox.showerror("错误", "未找到转换后的视频文件")
            return

        filename = filedialog.asksaveasfilename(
            defaultextension=".mp4",
            filetypes=[("MP4 files", "*.mp4")]
        )
        if filename:
            try:
                shutil.copy2(output_video, filename)
                messagebox.showinfo("成功", "视频已保存")
            except Exception as e:
                messagebox.showerror("错误", f"保存视频失败：{str(e)}")

    def load_config(self):
        try:
            if os.path.exists(self.config_path):
                with open(self.config_path, 'r', encoding='utf-8') as f:
                    config = json.load(f)
                    self.slide_duration.set(config.get('slide_duration', '5'))
                    self.transition_duration.set(config.get('transition_duration', '1'))
                    self.video_quality.set(config.get('video_quality', '高质量'))
                    self.bgm_volume.set(config.get('bgm_volume', '1.0'))
                    self.resolution.set(config.get('resolution', '自动'))
                    self.auto_next.set(config.get('auto_next', True))
                    self.save_text.set(config.get('save_text', True))
                    self.transition_effect.set(config.get('transition_effect', '无'))
        except Exception as e:
            messagebox.showwarning("警告", f"加载配置文件失败：{str(e)}")

    def save_config(self):
        try:
            config = {
                'slide_duration': self.slide_duration.get(),
                'transition_duration': self.transition_duration.get(),
                'video_quality': self.video_quality.get(),
                'bgm_volume': self.bgm_volume.get(),
                'resolution': self.resolution.get(),
                'auto_next': self.auto_next.get(),
                'save_text': self.save_text.get(),
                'transition_effect': self.transition_effect.get()
            }
            with open(self.config_path, 'w', encoding='utf-8') as f:
                json.dump(config, f, indent=4, ensure_ascii=False)
            messagebox.showinfo("成功", "配置已保存")
        except Exception as e:
            messagebox.showerror("错误", f"保存配置失败：{str(e)}")

    def get_ffmpeg_quality_params(self):
        quality_settings = {
            "低质量": ["-crf", "28", "-preset", "faster"],
            "中等质量": ["-crf", "23", "-preset", "medium"],
            "高质量": ["-crf", "18", "-preset", "slow"]
        }
        return quality_settings.get(self.video_quality.get(), quality_settings["高质量"])

    def select_ppt(self):
        filename = filedialog.askopenfilename(
            title="选择PPT文件",
            filetypes=[("PPT files", "*.pptx;*.ppt")]
        )
        if filename:
            self.ppt_path.set(filename)

    def select_bgm(self):
        filename = filedialog.askopenfilename(
            title="选择背景音乐",
            filetypes=[("Audio files", "*.mp3;*.wav;*.m4a")]
        )
        if filename:
            self.bgm_path.set(filename)

    def run(self):
        self.window.mainloop()


if __name__ == "__main__":
    app = PPTToVideo()
    app.run()